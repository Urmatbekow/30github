from pptx import Presentation
from pptx.util import Inches

def create_slide_deck(slides_content, output_file="presentation.pptx"):
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Your Presentation Title"
    subtitle.text = "Subtitle or Date"

    bullet_slide_layout = prs.slide_layouts[1]
    for slide_data in slides_content:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = slide_data['title']
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        if slide_data['content']:
            text_frame.text = slide_data['content'][0]
        for bullet in slide_data['content'][1:]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1

    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

slides_content = [
    {
        "title": "Introduction",
        "content": ["Overview of the topic", "Context", "Objectives"]
    },
    {
        "title": "Main Concepts",
        "content": ["Concept 1", "Concept 2", "Concept 3"]
    },
    {
        "title": "Conclusion",
        "content": ["Summary of key points", "Future work", "Q&A"]
    }
]

create_slide_deck(slides_content, "my_presentation.pptx")
